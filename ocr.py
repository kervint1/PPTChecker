import pytesseract
from pptx import Presentation
from PIL import Image
import io
import os
import sys

# PyInstallerがパスを正しく解決するための設定
if getattr(sys, 'frozen', False):
    # PyInstallerでビルドされた一つのファイルの実行可能ファイル内にいる場合
    base_path = sys._MEIPASS
else:
    # 通常のPythonインタプリタで実行されている場合
    base_path = os.path.abspath(".")

# Tesseractのパス設定
pytesseract.pytesseract.tesseract_cmd = os.path.join(base_path, 'tesseract', 'tesseract.exe')
os.environ['TESSDATA_PREFIX'] = os.path.join(base_path, 'tesseract', 'tessdata')

# def get_lp_account_name_message(file_path, left=55, top=10, right=200, bottom=100):
#     prs = Presentation(file_path)
#     slides_texts = []  # 各スライドのテキストリストを格納するリスト

#     for slide in prs.slides:
#         slide_texts = []  # 現在のスライドのテキストを格納するリスト
#         for shape in slide.shapes:
#             if shape.shape_type == 13 and shape.left.pt<37 and shape.left.pt>32 and shape.top.pt<147 and shape.top.pt>143:  # 画像タイプ
#                 image_stream = io.BytesIO(shape.image.blob)
#                 image = Image.open(image_stream)
#                 # 画像を指定された範囲でクロップ
#                 cropped_image = image.crop((left, top, right, bottom))
#                 text = pytesseract.image_to_string(cropped_image, lang='eng+jpn')
#                 slide_texts.append(text)
#         slides_texts.append(slide_texts)

#     return slides_texts

def get_lp_account_name_message(shape, left=85, top=10, right=210, bottom=50):
    shape_texts = None  # 現在のスライドのテキスト
    if shape.shape_type == 13:  # 画像タイプ
        image_stream = io.BytesIO(shape.image.blob)
        image = Image.open(image_stream)
        # 画像を指定された範囲でクロップ
        cropped_image = image.crop((left, top, right, bottom))
        text = pytesseract.image_to_string(cropped_image, lang='eng+jpn')
        shape_texts = text

        # # デバッグ用イメージ保存
        # temp_path = os.path.join(base_path, "debugpictures", "test_image.png")
        # cropped_image.save(temp_path)
    return shape_texts

# def get_lp_date_message(file_path, left=55, top=50, right=200, bottom=90, debug_dir=None):
#     prs = Presentation(file_path)
#     slides_texts = []  # 各スライドのテキストリストを格納するリスト
#     slide_number = 0  # スライドの番号を追跡

#     for slide in prs.slides:
#         slide_texts = []  # 現在のスライドのテキストを格納するリスト
#         shape_number = 0  # シェイプの番号を追跡

#         for shape in slide.shapes:
#             if shape.shape_type == 13 and 32 < shape.left.pt < 37 and 143 < shape.top.pt < 147:  # 画像タイプ
#                 image_stream = io.BytesIO(shape.image.blob)
#                 image = Image.open(image_stream)
#                 # 画像を指定された範囲でクロップ
#                 cropped_image = image.crop((left, top, right, bottom))
#                 # フォーマット指定で日付だけを読み取る設定
#                 custom_config = r'--psm 6 -c tessedit_char_whitelist=0123456789/ outputbase digits'
#                 text = pytesseract.image_to_string(cropped_image, lang='jpn', config=custom_config)
#                 slide_texts.append(text)

#                 # デバッグディレクトリが指定されている場合は、クロップした画像を保存
#                 if debug_dir:
#                     file_name = f"slide_{slide_number}_shape_{shape_number}.png"
#                     cropped_image.save(f"{debug_dir}/{file_name}")

#                 shape_number += 1
#         slides_texts.append(slide_texts)
#         slide_number += 1

#     return slides_texts

# # ファイルパスを指定して関数を呼び出し、結果を表示
# file_path1 = r"【事例資料】LOUIS VUITTON_LINE 公式アカウント_メッセージ配信_2024年1月以降.pptx"
# file_path2 = r"【事例資料】ヴァレンティノ_LINE 公式アカウント_メッセージ配信事例_2024年1月以降.pptx"
# file_path3 = r"【事例資料】ベイクルーズ_LINE 公式アカウント_メッセージ配信_2024年1月以降.pptx"
# # text_data = get_lp_date_message(file_path2, debug_dir=r"C:\Users\iniad\Documents\adpro\debugpictures")
# # print(text_data)
# from PIL import Image, ImageFilter, ImageEnhance
# import pytesseract

# # 画像のパス
# image_path = r"C:\Users\iniad\Documents\adpro\debugpictures\slide_0_shape_0.png"


# def preprocess_image(image_path):
#     # 画像を読み込む
#     image = Image.open(image_path)
    
#     # # コントラストを上げる
#     # enhancer = ImageEnhance.Contrast(image)
#     # image_enhanced = enhancer.enhance(2.0)  # コントラストの量を調整
    
#     # シャープネスを上げる
#     image_sharp = image.filter(ImageFilter.SHARPEN)
    
#     # # ノイズを減らす
#     # image_denoised = image_sharp.filter(ImageFilter.MedianFilter(size=3))
    
#     # # 画像をグレースケールに変換
#     # image_gray = image_denoised.convert('L')

#     # 一時ファイルとして保存し、OCRを実行する
#     temp_path = r"C:\Users\iniad\Documents\adpro\preprocessed.png"
#     image_sharp.save(temp_path)
#     text = pytesseract.image_to_string(Image.open(temp_path), lang='jpn')

#     return text, temp_path
# def enhance_sharpness(image_path, factor):
#     # 画像を読み込む
#     image = Image.open(image_path)
    
#     # シャープネスを上げる
#     sharpness_enhancer = ImageEnhance.Sharpness(image)
#     sharpened_image = sharpness_enhancer.enhance(factor)
    
#     # 画像を一時ファイルとして保存
#     temp_path = r"C:\Users\iniad\Documents\adpro\sharpened_image.png"
#     sharpened_image.save(temp_path)
#     text = pytesseract.image_to_string(Image.open(temp_path), lang='jpn+eng')
    
#     return text,temp_path

# # 前処理された画像からテキストを抽出する
# preprocessed_text, preprocessed_image_path = preprocess_image(image_path)
# preprocessed_text, preprocessed_image_path

# # 前処理された画像からテキストを抽出する
# # preprocessed_text = preprocess_image(r'C:\Users\iniad\Documents\adpro\debugpictures\slide_0_shape_0.png')
# # print(preprocessed_text)

# # sharpened_image_path = enhance_sharpness(preprocessed_image_path, 2.0)
# # print(sharpened_image_path)

# from pptx.enum.shapes import MSO_SHAPE

# def checkplace(filepath):
#     prs = Presentation(filepath)
#     slide_number = 0

#     for slide in prs.slides:
#         shapenumber = 0
#         slide_number+=1
#         for shape in slide.shapes:
#             shapenumber +=1
#             if shape.shape_type == 1:
#                 shapenumber +=1
#             else :
#                 print(slide_number,shapenumber,shape.top.pt,shape.left.pt)
#         if slide_number>50:
#             break
# def checkplace2(filepath,slidenumber):
#     prs = Presentation(filepath)
#     shapenumber = 0
#     for shape in prs.slides[slidenumber-1].shapes:
#         shapenumber +=1
#         print(slidenumber,shapenumber,shape.top.pt,shape.left.pt,shape.height.pt,shape.width.pt)

# import re
# def test1(filepath,slidenumber,shapenumber):
#     prs = Presentation(filepath)
#     shape = prs.slides[slidenumber-1].shapes[shapenumber-1]
    
#     return get_lp_account_name_message(shape)

# # print(test1(file_path3,145,1))

# # checkplace(file_path1)
# # checkplace2(file_path3,30)