from pptx import Presentation
import pandas as pd
from ocr import get_lp_account_name_message
from datetime import datetime
import os
import re

def extract_text_from_pptx_by_slide(file_path):
    slides_texts = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        slide_texts = []  # スライドごとのテキストを格納するリスト
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_texts.append(shape.text)
        slides_texts.append(slide_texts)
    return slides_texts

def classify_slide(slide,top_list,left_list):
    """
    スライドを分類する関数。表紙、月のスライド、内容のスライドのいずれかに分類。
    """
    objects = slide.shapes
    permissible = 5

    for shape in objects:
        if abs(shape.top.pt-top_list[0]) < permissible and abs(shape.left.pt-left_list[0]) < permissible:
            return 'cover'
        elif abs(shape.top.pt-top_list[1]) < permissible and abs(shape.left.pt-left_list[1]) < permissible:
            return 'month'
        elif abs(shape.top.pt-top_list[2]) < permissible and abs(shape.left.pt-left_list[2]) < permissible:
            return 'content'
    return None

def check_position(shape,permissible,cover_position_top,cover_position_left):
    if (abs(shape.top.pt-cover_position_top) < permissible and 
        abs(shape.left.pt-cover_position_left) < permissible):
        return True
    else :
        return False

def extract_cover_data(slide):
    """
    表紙のスライドからデータを抽出する関数。
    """
    objects = slide.shapes
    permissible = 20
    cover_position_top = [199,233]
    cover_position_left = [191,109]
    account_name = None
    message_or_voom = None
    error_message = None
    for shape in objects:
        if check_position(shape,permissible,cover_position_top[0],cover_position_left[0]):
            account_name = shape.text
        elif check_position(shape,permissible,cover_position_top[1],cover_position_left[1]):
            if re.search(r"メッセージ",shape.text):
                message_or_voom = 1
            elif re.search(r"VOOM",shape.text):
                message_or_voom = 2
        if account_name and message_or_voom:
            break
    if not(account_name and message_or_voom):
        error_message = "オブジェクトが基準値より20pt離れている"
    # print(account_name,message_or_voom,error_message)
    # 実際の実装はここに
    return pd.DataFrame([{
        'category_number': 1,
        'account_name': account_name,
        'message_or_voom': message_or_voom,        
        'error_message': error_message
    }])

def extract_month_data(slide):
    """
    月のスライドからデータを抽出する関数。
    """
    # 実際の実装はここに
    objects = slide.shapes
    permissible = 5
    cover_position_top = [199,233,283]
    cover_position_left = [191,109,111]
    account_name = None
    date = None
    error_message = None

    for shape in objects:
        if check_position(shape,permissible,cover_position_top[0],cover_position_left[0]):
            account_name = shape.text
        elif check_position(shape,permissible,cover_position_top[1],cover_position_left[1]):
            if re.search(r"メッセージ",shape.text):
                message_or_voom = 1
            elif re.search(r"VOOM",shape.text):
                message_or_voom = 2
        elif check_position(shape,permissible,cover_position_top[2],cover_position_left[2]):
            match = re.search(r"(\d{4})年(\d{1,2})月", shape.text)
            if match:
                year = int(match.group(1))
                month = int(match.group(2))
                date = datetime(year, month, 1, 0, 0, 0)
    if not(account_name and message_or_voom):
        error_message = "オブジェクトが基準値より20pt離れている"
    # print(account_name,message_or_voom,year,month,error_message)

    return pd.DataFrame([{
        'category_number': 2,
        'account_name': account_name,
        'message_or_voom': message_or_voom,  
        'date' : date,
        'error_message': error_message
    }])

def extract_content_data(slide,ocr):
    """
    内容のスライドからデータを抽出する関数。
    """
    # 実際の実装はここに
    
    objects = slide.shapes
    normal_permissible = 5
    strict_permissible = 1
    cover_position_top = [3,69]
    cover_position_left = [21,123]
    reference_line_top = [136,145,321,332]
    reference_line_left = [21,35,146,256,356,455,554,653]
    ad2_bottom = None
    ad1_width = 102.0455905511811
    lp_width = 93.54181102362205
    
    account_name = None
    message_or_voom = None
    date = None
    ad_presence = None
    ad_account_name = "ad_account"
    ad_number_count = 0
    lp_count = 0
    lp_number_count = 0
    arrow_presence = None
    error_message = ""
    numbercount = 0
    count_objects = 0


    for shape in objects:
        count_objects += 1
        #上部アカウント名VOOMまたはメッセージ
        if check_position(shape,normal_permissible,cover_position_top[0],cover_position_left[0]) and shape.shape_type != 13:  
            pattern = r'^([\w\s]+)　LINE公式アカウント\s+(\w+活用状況)\s*$'
            match = re.search(pattern, shape.text)
            if match:
                account_name = match.group(1).strip()
                message_voom = match.group(2).strip()
                if re.search(r"メッセージ",message_voom):
                    message_or_voom = 1
                elif re.search(r"VOOM",message_voom):
                    message_or_voom = 2
        #日付のボックス
        elif check_position(shape,normal_permissible,cover_position_top[1],cover_position_left[1]) and shape.shape_type != 13:
            # 正規表現パターンを定義
            pattern = r"(\d{1,2})月(\d{1,2})日.*?(\d{1,2}:\d{2})"
            
            # パターンにマッチするテキストを検索
            match = re.search(pattern, shape.text)
            
            if match:
                month = int(match.group(1))
                day = int(match.group(2))
                time = match.group(3)
                
                # 現在の年を取得
                year = datetime.now().year
                
                # datetime型に変換
                date_str = f"{year}-{month:02}-{day:02} {time}"
                date = datetime.strptime(date_str, "%Y-%m-%d %H:%M")
        #1.LPとADの大まかな範囲を指定(top厳密,left指定なし)
        #2.幅でpictureをLPとADを分類、位置判定、数
        #3.振り数を分類,位置判定,数
        #3.矢印有無
        elif (shape.top.pt>reference_line_top[0] - strict_permissible):
            if (shape.shape_type == 13 and abs(shape.width.pt - ad1_width) < 5):
                if check_position(shape,strict_permissible,reference_line_top[1],reference_line_left[1]):
                    ad_presence = True
                    if(ocr == True):
                        ad_account_name = get_lp_account_name_message(shape)
                    #145->ADのtop,371=>ADの最大height
                    if shape.top.pt +shape.height.pt > 145 + 371:
                        error_message += "adの高さが大きすぎる," + str(count_objects)
                elif check_position(shape,strict_permissible,reference_line_top[1],reference_line_left[2]):
                    ad2_bottom = shape.top.pt +shape.height.pt
                elif abs(shape.left.pt - reference_line_left[1]) < strict_permissible:
                    if shape.top.pt +shape.height.pt > 145 + 371:
                        error_message += "adの高さが大きすぎる," + str(count_objects)                      
                else:
                    error_message += "基準線に従っていないAD,"+str(count_objects)
            elif (shape.shape_type == 13 and abs(shape.width.pt - lp_width) < 5):
                lp_count += 1
            elif (shape.shape_type != 1 ):
                error_message +="基準線にあっていないオブジェクトがありますnot1。" + str(count_objects)
            elif (shape.auto_shape_type == 7):
                arrow_presence = shape.top.pt
            elif(shape.auto_shape_type == 1 and shape.left.pt > reference_line_left[3]-1):
                lp_number_count += 1
            elif(shape.auto_shape_type == 1 ):
                ad_number_count += 1
            else:
                error_message +="基準線にあっていないオブジェクトがあります1。,"
                

        # else:
            # error_message +="基準線にあっていないオブジェクトがあります2。,"
    
    if (ad2_bottom != None and arrow_presence):
        if(ad2_bottom>arrow_presence):
            error_message +="矢印がかぶっている"
    if not(account_name and message_or_voom):
        error_message += "オブジェクトが基準値より20pt離れている"
    # print(account_name,message_or_voom,month,day,time,ad_presence,ad_account_name,ad_number_count,lp_count,lp_number_count,arrow_presence,error_message)


    return pd.DataFrame([{
        'category_number': 3,
        'account_name': account_name,
        'message_or_voom': message_or_voom,
        'date': date,
        'ad_presence': ad_presence,
        'ad_account_name': ad_account_name,
        'ad_number_count': ad_number_count,
        'lp_count': lp_count,
        'lp_number_count': lp_number_count,
        'arrow_presence': arrow_presence,
        'error_message': error_message
    }])

# def summarize_slides(file_path):
#     """
#     スライドを分類し、それぞれのデータを取得してpandasでデータフレームにまとめる関数。
#     """
#     slides = Presentation(file_path).slides

#     #【事例資料】LOUIS VUITTON_LINE 公式アカウント_メッセージ配信_2024年1月以降.pptx(少数点以下切り捨て)
#     # 上記のpptでpositionを大まかに決める
#     # 1 4,2 6,3 6
#     LV_position_top = [474,233,3]
#     LV_position_left = [311,109,21]
#     permissible = 20
#     standard_top = []
#     standard_left = []
#     # その初期３スライドのpptの基準を決める
#     for i in range(0,3):
#         for shape in slides[i].shapes:
#             if (abs(shape.left.pt-LV_position_left[i])<permissible and
#                 abs(shape.top.pt-LV_position_top[i])<permissible):
#                 standard_top.append(round(shape.top.pt,0))
#                 standard_left.append(round(shape.left.pt))
#                 break
#             else:
#                 None
#     if len(standard_top)!=3:
#         print("top3 Slide Error")

#     # print(standard_top,standard_left)
#     data_frames = []

#     count= 0

#     for slide in slides:
#         count+=1
#         slide_type = classify_slide(slide,standard_top,standard_left)
#         if slide_type == 'cover':
#             df = extract_cover_data(slide)
#             # print(0)
#         elif slide_type == 'month':
#             df = extract_month_data(slide)
#             # print(1)
#         elif slide_type == 'content':
#             # print("content",count)
#             df = extract_content_data(slide)
            
#         else:
#             print(4)
#             df = pd.DataFrame([{
#                 'category_number': 4,
#                 'account_name': None,
#                 'error_message': 'No slide content'
#             }])
#         data_frames.append(df)
        
    
#     result_df = pd.concat(data_frames, ignore_index=True)
#     # print(result_df)
#     return result_df

def summarize_latest_slides(file_path, ocr, months = None,):
    """
    スライドを分類し、それぞれのデータを取得してpandasでデータフレームにまとめる関数。
    """
    slides = Presentation(file_path).slides

    #【事例資料】LOUIS VUITTON_LINE 公式アカウント_メッセージ配信_2024年1月以降.pptx(少数点以下切り捨て)
    # 上記のpptでpositionを大まかに決める
    # 1 4,2 6,3 6
    LV_position_top = [474,233,3]
    LV_position_left = [311,109,21]
    permissible = 20
    standard_top = []
    standard_left = []
    
    # その初期３スライドのpptの基準を決める
    for i in range(0, 3):
        for shape in slides[i].shapes:
            if (abs(shape.left.pt - LV_position_left[i]) < permissible and
                abs(shape.top.pt - LV_position_top[i]) < permissible):
                standard_top.append(round(shape.top.pt, 0))
                standard_left.append(round(shape.left.pt))
                break

    if len(standard_top) != 3:
        print("top3 Slide Error")
    
    data_frames = []
    count = 0
    month_count = 0
    
    for slide in reversed(slides):
        count += 1
        slide_type = classify_slide(slide, standard_top, standard_left)
        if slide_type == 'cover':
            df = extract_cover_data(slide)
        elif slide_type == 'month':
            df = extract_month_data(slide)
            month_count += 1
        elif slide_type == 'content':
            df = extract_content_data(slide,ocr)
        else:
            df = pd.DataFrame([{
                'category_number': 4,
                'account_name': None,
                'error_message': ['No slide content']
            }])
        
        data_frames.append(df)
        if months != None:
            if month_count == months:
                break

    data_frames.reverse()
    
    result_df = pd.concat(data_frames, ignore_index=True)
    return result_df


# ファイルパスを指定して関数を呼び出し、結果を表示します。
# print(extract_text_from_pptx_by_slide(file_path3))
# summarize_slides(file_path1).to_csv('file1Test1.csv')
# print(summarize_latest_slides(file_path1))
