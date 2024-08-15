# 必要なインポート
from tkinter import filedialog, Tk,messagebox
import os
from pptx import Presentation

# 定数の定義
reference_line_top = [136, 145, 321, 332]
reference_line_left = [21, 35, 146, 256, 356, 455, 554, 653]
strict_permissible = 1
normal_permissible = 10
huriban_width_height = 17.1

def change_font_to_yugothic():
    # ファイル選択
    root = Tk()
    root.withdraw()
    file_path = filedialog.askopenfilename()

    if os.path.exists(file_path):
        # ファイルを読み込む
        prs = Presentation(file_path)

        # フォントの変更
        for slide in prs.slides:
            for shape in slide.shapes:
                if (shape.has_text_frame and 
                    shape.top.pt > reference_line_top[0] - normal_permissible and
                    abs(shape.width.pt - huriban_width_height) + strict_permissible and
                    abs(shape.height.pt - huriban_width_height) + strict_permissible):
                    
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = '游ゴシック'


        # 変更を保存
        prs.save(file_path)
        messagebox.showinfo("変更結果","フォント変更が完了しました")

        print("フォント変更が完了しました")
    else:
        messagebox.showinfo("変更結果","ファイルが見つかりませんでした")
        print('ファイルが見つかりませんでした')

    root.mainloop()